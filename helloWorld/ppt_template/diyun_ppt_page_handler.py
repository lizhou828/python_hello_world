# -*- coding: utf-8 -*-
import os

from matplotlib.font_manager import FontProperties
from pptx.dml.color import RGBColor
from pptx.shapes.graphfrm import GraphicFrame
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.util import Cm  # Inches

import matplotlib.pyplot as plt
import matplotlib.ticker as mtick
import numpy as np
from helloWorld.ppt_template.screenshot_test import  getScreenShotForDiyunLandReport


class diyun_ppt_page_handler():

	@staticmethod
	def iter_cells(table):
		'''
		遍历table中的每一个单元格
		:param table:
		:return:
		'''
		for row in table.rows:
			for cell in row.cells:
				yield cell

	@staticmethod
	def handler_page_1(slide1):
		slide1_shapes = slide1.shapes
		# 第一页幻灯片处理
		title1 = slide1.shapes[5].text.format(landLocation="白云区白云新城AB2906009地块")
		slide1.shapes[5].text = ""
		p1 = slide1_shapes[5].text_frame.paragraphs[0]
		p1.alignment = PP_ALIGN.CENTER
		run1 = p1.add_run()
		run1.text = title1
		font1 = run1.font
		font1.name = '微软雅黑'
		font1.size = Pt(36)
		font1.bold = True
		font1.color.rgb = RGBColor(255, 255, 255)  # 白色字体

		title1 = slide1.shapes[3].text.format(year=2019, month=10)
		slide1.shapes[3].text = ""
		p1 = slide1_shapes[3].text_frame.paragraphs[0]
		p1.alignment = PP_ALIGN.LEFT
		run1 = p1.add_run()
		run1.text = title1
		font1 = run1.font
		font1.name = '微软雅黑'
		font1.size = Pt(18)
		font1.bold = True
		font1.color.rgb = RGBColor(0, 0, 0)  # 黑色字体
		# font1.underline = True  # 文字下划线
		# font1.italic = None


	@staticmethod
	def handler_page_4(slide4):
		slide4_shapes = slide4.shapes
		text = slide4.shapes[4].text.format(year=2018,city_name='广东省深圳市',year_gdp=13351.35,year_gdp_speed=6.5,per_person_gdp=86552,per_person_gdp_speed=7.6)
		slide4.shapes[4].text = ""
		p1 = slide4_shapes[4].text_frame.paragraphs[0]
		run1 = p1.add_run()
		run1.text = text

		# 数据的处理与转化
		table_data =[
						{"year":"2014","production_total":32230645,"growth_rate":6.5,"per_capita_production":"24000","per_capita_growth_rate":"1.4"},
					   	{"year": "2015", "production_total": 34649427, "growth_rate": 6, "per_capita_production": "25000","per_capita_growth_rate": "1.5"},
					    {"year": "2016", "production_total": 37155228, "growth_rate": 6.3, "per_capita_production": "26000","per_capita_growth_rate": "1.6"},
					    {"year": "2017", "production_total": 39784758, "growth_rate": 7, "per_capita_production": "27000","per_capita_growth_rate": "1.7"},
					    {"year": "2018", "production_total": 46177993, "growth_rate": 5, "per_capita_production": "28000","per_capita_growth_rate": "1.8"}
					 ]

		table_dict = {}
		year_list = []
		area_gdp_list =[]
		area_gdp_speed_list = []
		per_gdp_list = []
		per_gdp_speed_list = []
		for map in table_data:
			year = map.get("year")[2:4]
			table_dict["area_" + year] = map.get("production_total")
			table_dict["area_" + year + "s"] = map.get("growth_rate")
			table_dict["per_" + year] = map.get("per_capita_production")
			table_dict["per_" + year+ "s"] = map.get("per_capita_growth_rate")

			year_list.append(map.get("year"))
			area_gdp_list.append(map.get("production_total"))
			area_gdp_speed_list.append(map.get("growth_rate"))
			per_gdp_list.append(map.get("per_capita_production"))
			per_gdp_speed_list.append(map.get("per_capita_growth_rate"))


		# 把数据填充到table中
		graphicFrames = slide4.shapes
		for graphicFrame in graphicFrames:
			if type(graphicFrame) != GraphicFrame or not graphicFrame.has_table:
				continue
			for cell in diyun_ppt_page_handler.iter_cells(graphicFrame.table):
				if "{" not in cell.text and "}" not in cell.text:
					continue
				cell.text = cell.text.format(**table_dict)

			for cell in diyun_ppt_page_handler.iter_cells(graphicFrame.table):
				for paragraph in cell.text_frame.paragraphs:
					for run in paragraph.runs:
						run.font.size = Pt(12)
						run.font.color.rgb = RGBColor(0, 0, 0)  # 黑色字体


		#地区生产总值混合图
		l = [i for i in range(len(year_list))]  # 7组数据
		plt.rcParams['font.sans-serif'] = ['SimHei']  # 用来正常显示中文标签
		fmt = '%.2f%%'
		yticks = mtick.FormatStrFormatter(fmt)  # 设置百分比形式的坐标轴

		fig = plt.figure(figsize=(6,6),frameon=False)
		subplot1 = fig.add_subplot(111)


		subplot1.plot(l, area_gdp_speed_list, 'og-', label=u'增长速度')
		subplot1.yaxis.set_major_formatter(yticks)
		for i, (_x, _y) in enumerate(zip(l, area_gdp_speed_list)):
			plt.text(_x, _y, area_gdp_speed_list[i], color='black', fontsize=10, )  # 将数值显示在图形上
		subplot1.legend(loc=1)
		subplot1.set_ylim([min(area_gdp_speed_list), max(area_gdp_speed_list)])  # 设置y轴取值范围
		subplot1.set_ylabel(u'增长速度(%)')
		plt.legend(prop={'family': 'SimHei', 'size': 8})  # 右上角折线图标注，设置中文

		subplot1_twinx = subplot1.twinx()  # this is the important function
		plt.bar(l, area_gdp_list, width=0.2, alpha=0.5, color='#1E90FF', label=u'地区生产总值')  # 柱形图的数据、颜色、透明度等
		subplot1_twinx.legend(loc=2)
		subplot1_twinx.set_ylim([0, max(area_gdp_list)])  # 设置y轴取值范围
		subplot1_twinx.set_ylabel(u'地区生产总值(万元)')

		plt.legend(prop={'family': 'SimHei', 'size': 8}, loc="upper left")  # 左上角折线图标注，
		plt.xticks(l, year_list)
		plt.grid(axis='y', color='gray', linestyle=':', linewidth=1)
		slide4_chart1_path = "slide4_chart1.png"
		plt.savefig(slide4_chart1_path)
		left, top, width, height = Inches(0.5), Inches(1), Inches(4), Inches(4)
		slide4.shapes.add_picture(slide4_chart1_path, left, top, width, height)
		plt.cla()  # clean axis
		plt.clf()  # clean figure
		plt.close()#如果未另指定，则该窗口将是当前窗口。


		# 人均生产总值混合图
		plt.rcParams['font.sans-serif'] = ['SimHei']  # 用来正常显示中文标签
		fmt = '%.2f%%'
		yticks = mtick.FormatStrFormatter(fmt)  # 设置百分比形式的坐标轴

		fig = plt.figure(figsize=(6, 6), frameon=False)
		subplot1 = fig.add_subplot(111)

		subplot1.plot(l, per_gdp_speed_list, 'og-', label=u'增长速度')
		subplot1.yaxis.set_major_formatter(yticks)
		for i, (_x, _y) in enumerate(zip(l, per_gdp_speed_list)):
			plt.text(_x, _y, per_gdp_speed_list[i], color='black', fontsize=10, )  # 将数值显示在图形上
		subplot1.legend(loc=1)
		subplot1.set_ylim([min(per_gdp_speed_list), max(per_gdp_speed_list)])  # 设置y轴取值范围
		subplot1.set_ylabel(u'增长速度(%)')
		plt.legend(prop={'family': 'SimHei', 'size': 8})  # 右上角折线图标注，设置中文

		subplot1_twinx = subplot1.twinx()  # this is the important function
		plt.bar(l, per_gdp_list, width=0.2, alpha=0.5, color='#1E90FF', label=u'人均生产总值')  # 柱形图的数据、颜色、透明度等
		subplot1_twinx.legend(loc=2)
		subplot1_twinx.set_ylim([0, max(per_gdp_list)])  # 设置y轴取值范围
		subplot1_twinx.set_ylabel(u'人均生产总值(万元)')
		plt.grid(axis='y', color='gray', linestyle=':', linewidth=1)
		plt.legend(prop={'family': 'SimHei', 'size': 8}, loc="upper left")  # 左上角折线图标注，
		plt.xticks(l, year_list)
		slide4_chart1_path = "slide4_chart2.png"
		plt.savefig(slide4_chart1_path)
		left, top, width, height = Inches(5), Inches(1), Inches(4), Inches(4)
		slide4.shapes.add_picture(slide4_chart1_path, left, top, width, height)
		plt.cla()  # clean axis
		plt.clf()  # clean figure
		plt.close()  # 如果未另指定，则该窗口将是当前窗口。




	@staticmethod
	def handler_page_5(slide5):
		# 堆积柱形图：定义图表数据 ---------------------
		data_json=[
			{"year":"2013","first_industries":6504.0000,"second_industries":6657909.0000,"three_industries":83151073.0000},
			{"year":"2014", "first_industries": 5279.0000, "second_industries": 7824462.0000, "three_industries": 92194729.0000},
			{"year":"2015", "first_industries": 7584.0000, "second_industries": 7676087.0000, "three_industries": 103287649.0000},
			{"year":"2016", "first_industries": 8250.0000, "second_industries": 8315513.0000, "three_industries": 117607690.0000},
			{"year":"2017", "first_industries": 19551.0000, "second_industries": 9810979.0000, "three_industries": 131523886.0000},
		]
		year_list = []
		production1data=[]
		production2data = []
		production3data = []
		for data in data_json:
			year_list.append(data.get("year"))
			production1data.append(data.get("first_industries"))
			production2data.append(data.get("second_industries"))
			production3data.append(data.get("three_industries"))

		# 文本框中的文本替换
		all_production = production1data[-1] + production2data[-1] +production3data[-1]
		first_percent = '%.2f' % (production1data[-1] / all_production)
		first_percent = '%.2f' % (float(first_percent) * 100)
		second_percent = '%.2f' % (production2data[-1] / all_production)
		second_percent= '%.2f' % (float(second_percent) * 100)
		third_percent = '%.2f' % (production3data[-1] / all_production)
		third_percent = '%.2f' % (float(third_percent) * 100)
		slide6_shapes = slide5.shapes
		text = slide5.shapes[4].text.format(year=year_list[-1], first=production1data[-1], second=production2data[-1], third=production3data[-1],
											first_percent=first_percent,second_percent=second_percent,third_percent=third_percent)
		slide5.shapes[4].text = ""
		p1 = slide6_shapes[4].text_frame.paragraphs[0]
		run1 = p1.add_run()
		run1.text = text


		#堆积柱状图
		# 数据等位相加（np.array来实现）
		add_result_data1 = np.array(production1data) + np.array(production2data) + np.array(production3data)
		add_result_data2 = np.array(production1data) + np.array(production2data)
		add_result_data3 = np.array(production1data)
		plt.rcParams['font.sans-serif'] = ['SimHei']  # 用来正常显示中文标签
		plt.title("深圳市产业结构")
		plt.bar(year_list, add_result_data1, width =0.5, label=u'第三产业', fc='#FFD700')
		plt.bar(year_list, add_result_data2, width =0.5,label=u'第二产业', fc='#32CD32')
		plt.bar(year_list, add_result_data3, width =0.5,label=u'第一产业',tick_label=year_list, fc='#1E90FF')
		# 添加图例
		plt.ylabel(u'生产总值(万元)')
		plt.legend(loc="upper right")
		plt.grid(axis='y',color='gray',linestyle=':',linewidth=1)
		# # 根据数据显示合适的刻度
		# add_height =int( (max(add_result_data1) - min(add_result_data3))/10 )
		# plt.ylim(min(add_result_data3) - add_height ,max(add_result_data1) + add_height)

		pic_path = "3industries.png"
		plt.savefig(pic_path)
		# 堆积柱形图：将图表添加到幻灯片 --------------------
		left, top, width, height = Inches(1), Inches(1.4), Inches(8), Inches(3.5)
		slide5.shapes.add_picture(pic_path, left, top, width, height)
		plt.cla()  # clean axis
		plt.clf()  # clean figure
		plt.close()  # 如果未另指定，则该窗口将是当前窗口


		# 把数据填充到table中
		productionDataMap = {}
		productionDataMap["production0data"] = year_list
		productionDataMap["production1data"] = production1data
		productionDataMap["production2data"] = production2data
		productionDataMap["production3data"] = production3data
		graphicFrames = slide5.shapes
		for graphicFrame in graphicFrames:
			if type(graphicFrame) != GraphicFrame or not graphicFrame.has_table:
				continue
			for row_index in range(len(graphicFrame.table.rows)):
				row_data = productionDataMap.get("production{}data".format(row_index))
				for column_index in range(len(graphicFrame.table.rows[row_index].cells)):
					if column_index == 0:
						continue
					# 插入数据到 每个单元格
					graphicFrame.table.rows[row_index].cells[column_index].text = str(row_data[column_index-1])

					# 处理 每个单元格 的样式
					for paragraph in graphicFrame.table.rows[row_index].cells[column_index].text_frame.paragraphs:
						for run in paragraph.runs:
							run.font.size = Pt(12)
							run.font.color.rgb = RGBColor(0, 0, 0)  # 黑色字体


	@staticmethod
	def handler_page_6(slide6):
		data_json = [
			{"year": "2013", "city_name":"深圳市","permanentTotal": 1062.890, "registerTotal": 310.470, "density": 5323},
			{"year": "2014", "city_name":"深圳市","permanentTotal": 1077.890, "registerTotal": 332.210, "density": 5398},
			{"year": "2015", "city_name":"深圳市","permanentTotal": 1137.870, "registerTotal": 332.210, "density": 5697},
			{"year": "2016", "city_name":"深圳市","permanentTotal": 1190.840, "registerTotal": 384.520, "density": 5962},
			{"year": "2017", "city_name":"深圳市","permanentTotal": 1252.830, "registerTotal": 434.720, "density": 6234},
		]
		# 文本框中的文本替换
		slide6_shapes = slide6.shapes
		text = slide6.shapes[4].text.format(year=2018, city_name='广东省深圳市', register_total=434.720, permanent_total=1252.830,density=6234)
		slide6.shapes[4].text = ""
		p1 = slide6_shapes[4].text_frame.paragraphs[0]
		run1 = p1.add_run()
		run1.text = text

		# 数据格式转化
		year_list = []
		permanentTotal_list = []
		registerTotal_list = []
		density_list = []
		for data in data_json:
			year_list.append(data.get("year"))
			permanentTotal_list.append(data.get("permanentTotal"))
			registerTotal_list.append(data.get("registerTotal"))
			density_list.append(data.get("density"))


		# 并列柱状图与折线图
		l = [i for i in range(len(year_list))]  # n组数据
		plt.rcParams['font.sans-serif'] = ['SimHei']  # 用来正常显示中文标签
		fig = plt.figure(figsize=(6, 4))  # width, height in inches  ， If not provided, defaults to:rc:`figure.figsize` = ``[6.4, 4.8]``.
		ax1 = fig.add_subplot(111)
		plt.title("人口")
		ax1.plot(l, density_list, 'og-', label=u'人口密度')
		for i, (_x, _y) in enumerate(zip(l, density_list)):
			plt.text(_x, _y, density_list[i], color='black', fontsize=10, )  # 将数值显示在图形上
		ax1.legend(loc=1)
		ax1.set_ylim([min(density_list)-100, max(density_list)+100])  # 设置y轴取值范围
		ax1.set_ylabel(u'人口密度(人/平方千米)')
		plt.legend(prop={'family': 'SimHei', 'size': 8})  # 右上角折线图标注，设置中文

		# 柱状图，设置两组数与总宽度
		total_width, n = 0.8, 2
		width = total_width / n  # 每条柱状图的宽度

		ax2 = ax1.twinx()  # this is the important function
		plt.bar(l, permanentTotal_list, alpha=0.4, width=width, color='#FFD700', label=u'常住人口')  # 柱形图的数据、颜色、透明度等
		ax2.legend(loc=2)
		# ax2.set_ylim([min(permanentTotal_list)-100, max(permanentTotal_list)+100])  # 设置y轴取值范围
		ax2.set_ylabel(u'人口(万人)')
		for i in range(len(l)):
			l[i] = l[i] + width
		plt.bar(l, registerTotal_list, width=width, alpha=0.4, color='#1E90FF', label=u'户籍人口')
		plt.legend(prop={'family': 'SimHei', 'size': 8}, loc="upper left")  # 左上角折线图标注，
		plt.xticks(l, year_list)
		plt.grid(axis='y', color='gray', linestyle=':', linewidth=1)
		pic_path = "slide6_chart1.png"
		plt.savefig(pic_path)
		left, top, width, height = Inches(1.5), Inches(1), Inches(7), Inches(4)
		slide6.shapes.add_picture(pic_path, left, top, width, height)
		plt.cla()  # clean axis
		plt.clf()  # clean figure
		plt.close()  # 如果未另指定，则该窗口将是当前窗口

		# table中的数据填充
		# 把数据填充到table中
		rowDataMap = {}
		rowDataMap["row0data"] = year_list
		rowDataMap["row1data"] = permanentTotal_list
		rowDataMap["row2data"] = registerTotal_list
		rowDataMap["row3data"] = density_list
		graphicFrames = slide6.shapes
		for graphicFrame in graphicFrames:
			if type(graphicFrame) != GraphicFrame or not graphicFrame.has_table:
				continue
			for row_index in range(len(graphicFrame.table.rows)):
				row_data = rowDataMap.get("row{}data".format(row_index))
				for column_index in range(len(graphicFrame.table.rows[row_index].cells)):
					if column_index == 0:
						continue
					# 插入数据到 每个单元格
					graphicFrame.table.rows[row_index].cells[column_index].text = str(row_data[column_index - 1])

					# 处理 每个单元格 的样式
					for paragraph in graphicFrame.table.rows[row_index].cells[column_index].text_frame.paragraphs:
						for run in paragraph.runs:
							run.font.size = Pt(12)
							run.font.color.rgb = RGBColor(0, 0, 0)  # 黑色字体

	@staticmethod
	def handler_page_7(slide7):
		data_json = [
			{"year": "2013", "cityName": "深圳市", "disposableIncome": 44653.1000, "expenditure": 28812.4400},
			{"year": "2014", "cityName": "深圳市", "disposableIncome": 40948.0000, "expenditure": 28852.7100},
			{"year": "2015", "cityName": "深圳市", "disposableIncome": 44633.3000, "expenditure": 32359.2000},
			{"year": "2016", "cityName": "深圳市", "disposableIncome": 48695.0000, "expenditure": 36840.6100},
			{"year": "2017", "cityName": "深圳市", "disposableIncome": 52938.0000, "expenditure": 38320.1200},
		]
		# 数据格式转化
		year_list = []
		disposableIncome_list = []
		expenditure_list = []
		for data in data_json:
			year_list.append(data.get("year"))
			disposableIncome_list.append(data.get("disposableIncome"))
			expenditure_list.append(data.get("expenditure"))

		# 文本框内容替换
		incomePercent = '%.2f' % ((disposableIncome_list[-1] -  disposableIncome_list[-2]) / disposableIncome_list[-2])
		incomePercent = float(incomePercent) * 100
		if incomePercent == 0:
			incomePercent = "持平"
		elif incomePercent > 0:
			incomePercent = "增长" + str(incomePercent) + "%"
		else:
			incomePercent = "减少" + str(abs(incomePercent)) +"%"
		percent = '%.2f' % ( expenditure_list[-1] / disposableIncome_list[-1] )
		percent = float(percent) * 100
		slide7_shapes = slide7.shapes
		text = slide7.shapes[4].text.format(year=data_json[-1].get("year"), cityName=data_json[-1].get("cityName"),
											income = disposableIncome_list[-1],incomePercent=incomePercent,
											expenditure=expenditure_list[-1],percent = percent)
		slide7.shapes[4].text = ""
		p1 = slide7_shapes[4].text_frame.paragraphs[0]
		run1 = p1.add_run()
		run1.text = text


		# 并列柱状图
		x = list(range(len(year_list)))
		total_width, n = 0.8, 2
		width = total_width / n

		fig = plt.figure(figsize=(6, 4))  # width, height in inches  ， If not provided, defaults to:rc:`figure.figsize` = ``[6.4, 4.8]``.
		ax1 = fig.add_subplot(111)
		plt.title("居民收入")

		plt.bar(x, disposableIncome_list,width=width, color='#FFD700', label='人均可支配收入', tick_label=year_list)
		ax1.set_ylabel("单位：元")
		ax1.legend(loc=1)
		ax1.set_ylim(0, max(disposableIncome_list) + 100)  # 设置y轴取值范围
		plt.legend(prop={'family': 'SimHei', 'size': 8})  # 右上角折线图标注，设置中文

		for i in range(len(x)):
			x[i] = x[i] + width
		plt.bar(x, expenditure_list, width=width, color='#32CD32', label='人均消费性支出')
		plt.grid(axis='y', color='gray', linestyle=':', linewidth=1)
		plt.legend()
		pic_path = "slide7_chart1.png"
		plt.savefig(pic_path)
		left, top, width, height = Inches(1.5), Inches(1.3), Inches(6), Inches(3.8)
		slide7.shapes.add_picture(pic_path, left, top, width, height)
		plt.cla()  # clean axis
		plt.clf()  # clean figure
		plt.close()  # 如果未另指定，则该窗口将是当前窗口


		# table中的数据填充
		rowDataMap = {}
		rowDataMap["row0data"] = year_list
		rowDataMap["row1data"] = disposableIncome_list
		rowDataMap["row2data"] = expenditure_list
		graphicFrames = slide7.shapes
		for graphicFrame in graphicFrames:
			if type(graphicFrame) != GraphicFrame or not graphicFrame.has_table:
				continue
			for row_index in range(len(graphicFrame.table.rows)):
				row_data = rowDataMap.get("row{}data".format(row_index))
				for column_index in range(len(graphicFrame.table.rows[row_index].cells)):
					if column_index == 0:
						continue
					# 插入数据到 每个单元格
					cell_value = str(row_data[column_index - 1])
					graphicFrame.table.rows[row_index].cells[column_index].text = cell_value

					# 处理 每个单元格 的样式
					for paragraph in graphicFrame.table.rows[row_index].cells[column_index].text_frame.paragraphs:
						for run in paragraph.runs:
							run.font.size = Pt(12)
							run.font.color.rgb = RGBColor(0, 0, 0)  # 黑色字体

	@staticmethod
	def handler_page_8(slide8):
		data_json = [
			{"year": "2013", "cityName":"深圳市","fixedAssets": 24901970.000, "realEstate": 876.900},
			{"year": "2014", "cityName":"深圳市","fixedAssets": 27174226.000, "realEstate": 1069.490},
			{"year": "2015", "cityName":"深圳市","fixedAssets": 32983076.000, "realEstate": 1331.030},
			{"year": "2016", "cityName":"深圳市","fixedAssets": 40781638.000, "realEstate": 1756.520},
			{"year": "2017", "cityName":"深圳市","fixedAssets": 51473152.000, "realEstate": 2135.860},
		]

		# 数据格式转化
		year_list = []
		fixedAssets_list = []
		realEstate_list = []
		percent_list = []
		for data in data_json:
			year_list.append(data.get("year"))
			fixedAssets = data.get("fixedAssets") /10000
			fixedAssets_list.append(fixedAssets)
			realEstate_list.append(data.get("realEstate"))
			percent ='%.2f' %(data.get("realEstate") / fixedAssets) #保留2位小数
			percent_list.append(float(percent) * 100)

		# 文本框中的文本替换
		yearPercent = '%.2f' % ((fixedAssets_list[-1] -  fixedAssets_list[-2]) / fixedAssets_list[-2])
		yearPercent = float(yearPercent) * 100
		if yearPercent == 0:
			yearPercent = "持平"
		elif yearPercent > 0:
			yearPercent = "增长" + str(yearPercent) + "%"
		else:
			yearPercent = "减少" + str(abs(yearPercent)) +"%"

		slide6_shapes = slide8.shapes
		text = slide8.shapes[4].text.format(year=data_json[-1].get("year"), cityName=data_json[-1].get("cityName"),
											fixedAssets=fixedAssets_list[-1], realEstate=realEstate_list[-1],
											realEstatePercent=percent_list[-1],yearPercent=yearPercent)
		slide8.shapes[4].text = ""
		p1 = slide6_shapes[4].text_frame.paragraphs[0]
		run1 = p1.add_run()
		run1.text = text


		# 并列柱状图与折线图
		l = [i for i in range(len(year_list))]  # n组数据
		plt.rcParams['font.sans-serif'] = ['SimHei']  # 用来正常显示中文标签
		fig = plt.figure(figsize=(6, 4))  # width, height in inches  ， If not provided, defaults to:rc:`figure.figsize` = ``[6.4, 4.8]``.
		ax1 = fig.add_subplot(111)
		plt.title("资产投资")
		ax1.plot(l, percent_list, 'og-', label=u'房投/固投')
		fmt = '%.2f%%'
		yticks = mtick.FormatStrFormatter(fmt)  # 设置百分比形式的坐标轴
		ax1.yaxis.set_major_formatter(yticks)
		for i, (_x, _y) in enumerate(zip(l, percent_list)):
			plt.text(_x, _y, str(percent_list[i])+"%", color='black', fontsize=10, )  # 将数值显示在图形上
		ax1.legend(loc=1)
		ax1.set_ylim(0,max(percent_list)+10)  # 设置y轴取值范围
		ax1.set_ylabel(u'房投/固投(百分比)')
		plt.legend(prop={'family': 'SimHei', 'size': 8})  # 右上角折线图标注，设置中文

		# 柱状图，设置两组数与总宽度
		total_width, n = 0.8, 2
		width = total_width / n  # 每条柱状图的宽度

		ax2 = ax1.twinx()  # this is the important function
		plt.bar(l, fixedAssets_list, alpha=0.4, width=width, color='#FFD700', label=u'固定资产投资')  # 柱形图的数据、颜色、透明度等
		ax2.legend(loc=2)
		# ax2.set_ylim([min(permanentTotal_list)-100, max(permanentTotal_list)+100])  # 设置y轴取值范围
		ax2.set_ylabel(u'投资(亿元)')
		for i in range(len(l)):
			l[i] = l[i] + width
		plt.bar(l, realEstate_list, width=width, alpha=0.4, color='#1E90FF', label=u'房地产投资')
		plt.legend(prop={'family': 'SimHei', 'size': 8}, loc="upper left")  # 左上角折线图标注，
		plt.xticks(l, year_list)
		plt.grid(axis='y', color='gray', linestyle=':', linewidth=1)
		pic_path = "slide8_chart1.png"
		plt.savefig(pic_path)
		left, top, width, height = Inches(1.5), Inches(1.2), Inches(6), Inches(3.8)
		slide8.shapes.add_picture(pic_path, left, top, width, height)
		plt.cla()  # clean axis
		plt.clf()  # clean figure
		plt.close()  # 如果未另指定，则该窗口将是当前窗口

		# table中的数据填充
		rowDataMap = {}
		rowDataMap["row0data"] = year_list
		rowDataMap["row1data"] = fixedAssets_list
		rowDataMap["row2data"] = realEstate_list
		rowDataMap["row3data"] = percent_list
		graphicFrames = slide8.shapes
		for graphicFrame in graphicFrames:
			if type(graphicFrame) != GraphicFrame or not graphicFrame.has_table:
				continue
			for row_index in range(len(graphicFrame.table.rows)):
				row_data = rowDataMap.get("row{}data".format(row_index))
				for column_index in range(len(graphicFrame.table.rows[row_index].cells)):
					if column_index == 0:
						continue
					# 插入数据到 每个单元格
					if row_index == 3 :
						cell_value = str(row_data[column_index - 1]) + "%"
					else:
						cell_value = str(row_data[column_index - 1])
					graphicFrame.table.rows[row_index].cells[column_index].text = cell_value

					# 处理 每个单元格 的样式
					for paragraph in graphicFrame.table.rows[row_index].cells[column_index].text_frame.paragraphs:
						for run in paragraph.runs:
							run.font.size = Pt(12)
							run.font.color.rgb = RGBColor(0, 0, 0)  # 黑色字体

	@staticmethod
	def handler_page_9(slide9):
		data_json = [
			{"year": "2014", "cityName": "深圳市", "supplyArea": 1107985.0000, "dealArea": 167985.0000 ,"dealFloorPrice":5500.0000},
			{"year": "2015", "cityName": "深圳市", "supplyArea": 1108985.0000, "dealArea": 165985.0000,"dealFloorPrice":15000.0000},
			{"year": "2016", "cityName": "深圳市", "supplyArea": 142985.0000, "dealArea": 15985.0000,"dealFloorPrice":55000.0000},
			{"year": "2017", "cityName": "深圳市", "supplyArea": 152985.0000, "dealArea": 19985.0000,"dealFloorPrice":56000.0000},
			{"year": "2018", "cityName": "深圳市", "supplyArea": 152585.0000, "dealArea": 19585.0000,"dealFloorPrice":57000.0000},
		]

		# 数据格式转化
		year_list = []
		supplyArea_list = []
		dealArea_list = []
		dealFloorPrice_list = []
		for data in data_json:
			year_list.append(data.get("year"))
			supplyArea_list.append(data.get("supplyArea"))
			dealArea_list.append(data.get("dealArea"))
			dealFloorPrice_list.append(data.get("dealFloorPrice"))


		# 文本框中的文本替换
		supplyPercent = '%.2f' % ((supplyArea_list[-1] - supplyArea_list[-2]) / supplyArea_list[-2])
		supplyPercent = float(supplyPercent) * 100
		if supplyPercent == 0:
			supplyPercent = "持平"
		elif supplyPercent > 0:
			supplyPercent = "增长" + str(supplyPercent) + "%"
		else:
			supplyPercent = "减少" + str(abs(supplyPercent)) +"%"

		dealAreaPercent = '%.2f' % ((dealArea_list[-1] - dealArea_list[-2]) / dealArea_list[-2])
		dealAreaPercent = float(dealAreaPercent) * 100
		if dealAreaPercent == 0:
			dealAreaPercent = "持平"
		elif dealAreaPercent > 0:
			dealAreaPercent = "增长" + str(dealAreaPercent) + "%"
		else:
			dealAreaPercent = "减少" + str(abs(dealAreaPercent)) + "%"

		slide6_shapes = slide9.shapes
		text = slide9.shapes[4].text.format(year=data_json[-1].get("year"), cityName=data_json[-1].get("cityName"),
											supplyArea=supplyArea_list[-1], dealArea=dealArea_list[-1],dealFloorPrice=dealFloorPrice_list[-1],
											supplyPercent=supplyPercent,dealAreaPercent=dealAreaPercent)
		slide9.shapes[4].text = ""
		p1 = slide6_shapes[4].text_frame.paragraphs[0]
		run1 = p1.add_run()
		run1.text = text




		# 并列柱状图与折线图
		l = [i for i in range(len(year_list))]  # n组数据
		plt.rcParams['font.sans-serif'] = ['SimHei']  # 用来正常显示中文标签
		fig = plt.figure(figsize=(8, 3.5))  # width, height in inches  ， If not provided, defaults to:rc:`figure.figsize` = ``[6.4, 4.8]``.
		ax1 = fig.add_subplot(111)
		plt.title("土地市场成交走势")
		ax1.plot(l, dealFloorPrice_list, 'og-', label=u'成交楼面价')
		for i, (_x, _y) in enumerate(zip(l, dealFloorPrice_list)):
			plt.text(_x, _y, dealFloorPrice_list[i], color='black', fontsize=10, )  # 将数值显示在图形上
		ax1.legend(loc=1)
		ax1.set_ylim(0, max(dealFloorPrice_list) )  # 设置y轴取值范围
		ax1.set_ylabel(u'成交楼面价(元/平米)')
		plt.legend(prop={'family': 'SimHei', 'size': 8})  # 右上角折线图标注，设置中文

		# 柱状图，设置两组数与总宽度
		total_width, n = 0.8, 2
		width = total_width / n  # 每条柱状图的宽度

		ax2 = ax1.twinx()  # this is the important function
		plt.bar(l, supplyArea_list, alpha=0.4, width=width, color='#FFD700', label=u'供应面积')  # 柱形图的数据、颜色、透明度等
		ax2.legend(loc=2)
		# ax2.set_ylim( min(supplyArea_list), max(supplyArea_list) + 1000)  # 设置y轴取值范围
		ax2.set_ylabel(u'面积(万平方米)')
		for i in range(len(l)):
			l[i] = l[i] + width
		plt.bar(l, dealArea_list, width=width, alpha=0.4, color='#1E90FF', label=u'成交面积')
		plt.legend(prop={'family': 'SimHei', 'size': 8}, loc="upper left")  # 左上角折线图标注，
		plt.xticks(l, year_list)
		plt.grid(axis='y', color='gray', linestyle=':', linewidth=1)
		pic_path = "slide9_chart1.png"
		plt.savefig(pic_path)
		left, top, width, height = Inches(1.5), Inches(1.2), Inches(6), Inches(3.8)
		slide9.shapes.add_picture(pic_path, left, top, width, height)
		plt.cla()  # clean axis
		plt.clf()  # clean figure
		plt.close()  # 如果未另指定，则该窗口将是当前窗口

		# table中的数据填充
		rowDataMap = {}
		rowDataMap["row0data"] = year_list
		rowDataMap["row1data"] = supplyArea_list
		rowDataMap["row2data"] = dealArea_list
		rowDataMap["row3data"] = dealFloorPrice_list
		graphicFrames = slide9.shapes
		for graphicFrame in graphicFrames:
			if type(graphicFrame) != GraphicFrame or not graphicFrame.has_table:
				continue
			for row_index in range(len(graphicFrame.table.rows)):
				row_data = rowDataMap.get("row{}data".format(row_index))
				for column_index in range(len(graphicFrame.table.rows[row_index].cells)):
					if column_index == 0:
						continue
					# 插入数据到 每个单元格
					cell_value = str(row_data[column_index - 1])
					graphicFrame.table.rows[row_index].cells[column_index].text = cell_value

					# 处理 每个单元格 的样式
					for paragraph in graphicFrame.table.rows[row_index].cells[column_index].text_frame.paragraphs:
						for run in paragraph.runs:
							run.font.size = Pt(12)
							run.font.color.rgb = RGBColor(0, 0, 0)  # 黑色字体

	@staticmethod
	def handler_page_10(slide10):
		house_sale_json = [
			{"year": "2013", "cityName": "深圳市", "saleArea":527.1600, "saleAveragePrice":21808.0000 },
			{"year": "2014", "cityName": "深圳市", "saleArea":474.8100, "saleAveragePrice":23955.0000},
			{"year": "2015", "cityName": "深圳市", "saleArea":747.8300, "saleAveragePrice":33406.0000},
			{"year": "2016", "cityName": "深圳市", "saleArea":660.0800, "saleAveragePrice":53455.0000},
			{"year": "2017", "cityName": "深圳市", "saleArea":520.9700, "saleAveragePrice":54455.0000}
		]
		house_dev_json = [
			{"year": "2013", "cityName": "深圳市", "developmentArea": 4003.4900, "completedArea": 353.5500},
			{"year": "2014", "cityName": "深圳市", "developmentArea": 4492.1800, "completedArea": 425.3100},
			{"year": "2015", "cityName": "深圳市", "developmentArea": 4978.4100, "completedArea": 360.2100},
			{"year": "2016", "cityName": "深圳市", "developmentArea": 5173.9900, "completedArea": 490.0300},
			{"year": "2017", "cityName": "深圳市", "developmentArea": 5709.3400, "completedArea": 285.0600}
		]
		# 数据格式转化
		dev_year_list = []
		developmentArea_list = []
		completedArea_list = []
		for data in house_dev_json:
			dev_year_list.append(data.get("year"))
			developmentArea_list.append(data.get("developmentArea"))
			completedArea_list.append(data.get("completedArea"))

		sale_year_list = []
		saleArea_list = []
		saleAveragePrice_list = []
		for data in house_sale_json:
			sale_year_list.append(data.get("year"))
			saleArea_list.append(data.get("saleArea"))
			saleAveragePrice_list.append(data.get("saleAveragePrice"))


		# 文本框中的文本替换
		saleAreaPercent = '%.2f' % ( (saleArea_list[-1] - saleArea_list[-2]) / saleArea_list[-2] )
		saleAreaPercent = float(saleAreaPercent) * 100
		if saleAreaPercent == 0:
			saleAreaPercent = "持平"
		elif saleAreaPercent > 0:
			saleAreaPercent = "增长" + str(saleAreaPercent) + "%"
		else:
			saleAreaPercent = "减少" + str(abs(saleAreaPercent)) + "%"
		#文本框的内容替换
		slide10_shapes = slide10.shapes
		text = slide10.shapes[3].text.format(year=sale_year_list[-1], cityName=house_sale_json[-1].get("cityName"),
											saleArea=saleArea_list[-1], saleAveragePrice=saleAveragePrice_list[-1],
											 developmentArea=developmentArea_list[-1],completedArea=completedArea_list[-1],
											 saleAreaPercent=saleAreaPercent)
		slide10.shapes[3].text = ""
		p1 = slide10_shapes[3].text_frame.paragraphs[0]
		run1 = p1.add_run()
		run1.text = text



		# 第一个并列柱状图
		l = [i for i in range(len(sale_year_list))]  # 7组数据
		plt.rcParams['font.sans-serif'] = ['SimHei']  # 用来正常显示中文标签
		fig = plt.figure(figsize=(5.2, 4), frameon=False)
		subplot1 = fig.add_subplot(111)
		plt.title("商品住宅销售情况")
		subplot1.plot(l, saleAveragePrice_list, 'og-', label=u'销售均价')
		for i, (_x, _y) in enumerate(zip(l, saleAveragePrice_list)):
			plt.text(_x, _y, saleAveragePrice_list[i], color='black', fontsize=10, )  # 将数值显示在图形上
		subplot1.legend(loc=1)
		subplot1.set_ylim([min(saleAveragePrice_list), max(saleAveragePrice_list)])  # 设置y轴取值范围
		subplot1.set_ylabel(u'销售均价(元/平米)')
		plt.legend(prop={'family': 'SimHei', 'size': 8})  # 右上角折线图标注，设置中文

		subplot1_twinx = subplot1.twinx()  # this is the important function
		plt.bar(l, saleArea_list, width=0.2, alpha=0.5, color='#1E90FF', label=u'销售面积')  # 柱形图的数据、颜色、透明度等
		subplot1_twinx.legend(loc=2)
		subplot1_twinx.set_ylim([0, max(saleArea_list)])  # 设置y轴取值范围
		subplot1_twinx.set_ylabel(u'销售面积(万平方米)')

		plt.legend(prop={'family': 'SimHei', 'size': 8}, loc="upper right")  # 左上角折线图标注，
		plt.xticks(l, sale_year_list)
		plt.grid(axis='y', color='gray', linestyle=':', linewidth=1)
		slide10_chart1_path = "slide10_chart1.png"
		plt.savefig(slide10_chart1_path)
		left, top, width, height = Inches(0.1), Inches(1.3), Inches(5), Inches(4)
		slide10.shapes.add_picture(slide10_chart1_path, left, top, width, height)
		plt.cla()  # clean axis
		plt.clf()  # clean figure
		plt.close()  # 如果未另指定，则该窗口将是当前窗口。



		# 第二个并列柱状图
		x = list(range(len(dev_year_list)))
		total_width, n = 0.4, 2
		width = total_width / n

		fig = plt.figure(figsize=(4.8, 4))  # width, height in inches  ， If not provided, defaults to:rc:`figure.figsize` = ``[6.4, 4.8]``.
		ax1 = fig.add_subplot(111)
		plt.title("商品住宅开发情况")

		plt.bar(x, developmentArea_list, width=width, color='#FFD700', label='施工面积', tick_label=dev_year_list)
		ax1.set_ylabel("单位：万平米")
		ax1.legend(loc=1)
		# ax1.set_ylim(0, max(developmentArea_list) + 100)  # 设置y轴取值范围
		plt.legend(prop={'family': 'SimHei', 'size': 8})  # 右上角折线图标注，设置中文

		for i in range(len(x)):
			x[i] = x[i] + width
		plt.bar(x, completedArea_list, width=width, color='#32CD32', label='竣工面积')
		plt.grid(axis='y', color='gray', linestyle=':', linewidth=1)
		plt.legend()
		pic_path = "slide10_chart2.png"
		plt.savefig(pic_path)
		left, top, width, height = Inches(5.4), Inches(1.3), Inches(5), Inches(4)
		slide10.shapes.add_picture(pic_path, left, top, width, height)
		plt.cla()  # clean axis
		plt.clf()  # clean figure
		plt.close()  # 如果未另指定，则该窗口将是当前窗口

		# table中的数据填充
		table1RowDataMap = {}
		table1RowDataMap["row0data"] = sale_year_list
		table1RowDataMap["row1data"] = saleArea_list
		table1RowDataMap["row2data"] = saleAveragePrice_list

		table2RowDataMap = {}
		table2RowDataMap["row0data"] = dev_year_list
		table2RowDataMap["row1data"] = developmentArea_list
		table2RowDataMap["row2data"] = completedArea_list

		graphicFrames = slide10.shapes
		frame_index = 0
		for graphicFrame in graphicFrames:
			frame_index +=1
			if type(graphicFrame) != GraphicFrame or not graphicFrame.has_table:
				continue
			if frame_index != 5 and frame_index != 6:
				continue
			for row_index in range(len(graphicFrame.table.rows)):
				row_data = []
				if frame_index == 5:
					row_data = table1RowDataMap.get("row{}data".format(row_index))
				if frame_index == 6:
					row_data = table2RowDataMap.get("row{}data".format(row_index))

				for column_index in range(len(graphicFrame.table.rows[row_index].cells)):
					if column_index == 0:
						continue
					# 插入数据到 每个单元格
					cell_value = str(row_data[column_index - 1])
					graphicFrame.table.rows[row_index].cells[column_index].text = cell_value

					# 处理 每个单元格 的样式
					for paragraph in graphicFrame.table.rows[row_index].cells[column_index].text_frame.paragraphs:
						for run in paragraph.runs:
							run.font.size = Pt(12)
							run.font.color.rgb = RGBColor(0, 0, 0)  # 黑色字体


	@staticmethod
	def handler_page_12(slide12):
		# url = "file:///" + os.getcwd() + os.sep + "baidu_maps_in_ppt.html"
		# pic_path = getScreenShotForDiyunLandReport(url)

		pic_path = "baidu_maps_in_ppt1.png"  # 指定保存文件的文件名

		top = Inches(0.7)
		left = Inches(0.1)
		width = Inches(6)
		height = Inches(5)
		slide12.shapes.add_picture(pic_path, left, top, width, height)


	@staticmethod
	def handler_page_14(slide14):
		data_dict = {'land_location':'白云区白云新城AB2906009地块',"land_sn":"ab2906009","land_use_details":"城镇住宅用地","type":"挂牌","land_total_area":"67695.0000",
					 "sale_time":"70","plot_ratio":"大于1并且小于或等于5.5","building_density":"小于或等于30","greening_rate":"大于或等于35","building_limited_height":"小于或等于120",
					 "open_end_time":"2019-06-14 10:00:00","cash_deposit":"12300.0000","publish_time":"2019-05-16 00:00:00","open_start_time":"2019-06-04 00:00:00","starting_price":"193660.0000","price_increase":"3000.0000",
					 "contract_signing_date":"2019-06-25 00:00:00","transaction_price":"280807.0000","land_owner":"广州市瑞业房地产开发有限公司","floor_area_price":"","premium_rate":"0.13"}

		graphicFrames = slide14.shapes
		for graphicFrame in graphicFrames:
			if type(graphicFrame) != GraphicFrame or not graphicFrame.has_table:
				continue
			for cell in diyun_ppt_page_handler.iter_cells(graphicFrame.table):
				if "{" not in cell.text and "}" not in cell.text:
					continue
				cell.text = cell.text.format(**data_dict)
				# for paragraph in cell.text_frame.paragraphs:
				# 	for run in paragraph.runs:
				# 		run.text = run.text.format(data_dict)
						# run.font.size = Pt(12)
						# run.font.color.rgb = RGBColor(0, 0, 0)  # 黑色字体

			for cell in diyun_ppt_page_handler.iter_cells(graphicFrame.table):
				for paragraph in cell.text_frame.paragraphs:
					for run in paragraph.runs:
						run.font.size = Pt(12)
						run.font.color.rgb = RGBColor(0, 0, 0)  # 黑色字体

	@staticmethod
	def handler_page_16(slide16):
		land_datas = []
		land_data1 = {'land_location': '白云区白云新城AB2906009地块', "land_sn": "ab2906009", "land_use_details": "城镇住宅用地", "type": "挂牌",
		 "land_total_area": "67695.0000",
		 "sale_time": "70", "plot_ratio": "大于1并且小于或等于5.5", "building_density": "小于或等于30", "greening_rate": "大于或等于35",
		 "building_limited_height": "小于或等于120",
		 "open_end_time": "2019-06-14 10:00:00", "cash_deposit": "12300.0000", "publish_time": "2019-05-16 00:00:00",
		 "open_start_time": "2019-06-04 00:00:00", "starting_price": "193660.0000", "price_increase": "3000.0000",
		 "contract_signing_date": "2019-06-25 00:00:00", "transaction_price": "280807.0000",
		 "land_owner": "广州市瑞业房地产开发有限公司", "floor_area_price": "", "premium_rate": "0.13","lng": "113.1051556943977", "lat": "23.0192965357819"}
		land_datas.append(land_data1)

	@staticmethod
	def handler_page_18(slide18):
		table_data = []
		row_data1 = {"name":"东方盛世花园二期","time":"2018-10-08","type":"板楼","price":"73638.73","lng":"113.1051556943977","lat":"23.0192965357819"}
		table_data.append(row_data1)
		row_data2 = {"name":"深业泰富广场","time":"2019-08-08","type":"板塔结合","price":"59228.69","lng":"113.1051356943977","lat":"23.0191965357819"}
		table_data.append(row_data2)
		row_data3 = {"name": "深业泰富广场", "time": "2019-08-08", "type": "板塔结合", "price": "59228.69", "lng": "113.1051356943977", "lat": "23.0191965357819"}
		table_data.append(row_data3)
		row_data4 = {"name": "深业泰富广场", "time": "2019-08-08", "type": "板塔结合", "price": "59228.69","lng": "113.1051356943977", "lat": "23.0191965357819"}
		table_data.append(row_data4)
		row_data5 = {"name": "深业泰富广场", "time": "2019-08-08", "type": "板塔结合", "price": "59228.69","lng": "113.1051356943977", "lat": "23.0191965357819"}
		table_data.append(row_data5)

		row_data6 = {"name": "深业泰富广场", "time": "2019-08-08", "type": "板塔结合", "price": "59228.69","lng": "113.1051356943977", "lat": "23.0191965357819"}
		table_data.append(row_data6)
		row_data7 = {"name": "深业泰富广场", "time": "2019-08-08", "type": "板塔结合", "price": "59228.69","lng": "113.1051356943977", "lat": "23.0191965357819"}
		table_data.append(row_data7)
		row_data8 = {"name": "深业泰富广场", "time": "2019-08-08", "type": "板塔结合", "price": "59228.69","lng": "113.1051356943977", "lat": "23.0191965357819"}
		table_data.append(row_data8)
		row_data9 = {"name": "深业泰富广场", "time": "2019-08-08", "type": "板塔结合", "price": "59228.69","lng": "113.1051356943977", "lat": "23.0191965357819"}
		table_data.append(row_data9)
		row_data10 = {"name": "深业泰富广场", "time": "2019-08-08", "type": "板塔结合", "price": "59228.69","lng": "113.1051356943977", "lat": "23.0191965357819"}
		table_data.append(row_data10)

		# 表格样式 --------------------
		rows = len(table_data)+1
		cols = 5
		top = Cm(2)
		left = Cm(13.5)  # Inches(2.0)
		width = Cm(10)  # Inches(6.0)
		height = Cm(2)  # Inches(0.8)

		# 添加表格到幻灯片 --------------------
		table = slide18.shapes.add_table(rows, cols, left, top, width, height).table
		# 设置单元格宽度
		table.columns[0].width = Cm(1.5)  # Inches(2.0)
		table.columns[1].width = Cm(3)
		table.columns[2].width = Cm(3)
		table.columns[3].width = Cm(2)
		table.columns[4].width = Cm(2.5)

		# 插入数据到table
		for row_index in range(len(table.rows)):
			if row_index == 0:
				table.rows[row_index].cells[0].text="编号"
				table.rows[row_index].cells[1].text = "名称"
				table.rows[row_index].cells[2].text = "开盘时间"
				table.rows[row_index].cells[3].text = "产品类型"
				table.rows[row_index].cells[4].text = "成交价格(元/㎡)"
				continue
			for cell_index in range(len(table.rows[row_index].cells)):
				if cell_index == 0:
					table.rows[row_index].cells[cell_index].text = str(row_index)
				if cell_index == 1:
					table.rows[row_index].cells[cell_index].text = table_data[row_index-1].get("name")
				if cell_index == 2:
					table.rows[row_index].cells[cell_index].text = table_data[row_index-1].get("time")
				if cell_index == 3:
					table.rows[row_index].cells[cell_index].text = table_data[row_index-1].get("type")
				if cell_index == 4:
					table.rows[row_index].cells[cell_index].text = table_data[row_index-1].get("price")

		# 处理table的样式、字体等
		for cell in diyun_ppt_page_handler.iter_cells(table):
					for paragraph in cell.text_frame.paragraphs:
						for run in paragraph.runs:
							run.font.size = Pt(12)
							run.font.color.rgb = RGBColor(0, 0, 0)  # 黑色字体

	@staticmethod
	def handler_page_20(slide20):
		table_data = []
		row_data1 = {"name":"宝安国际机场（飞机场）","distance":"0.8","lng":"113.1051556943977","lat":"23.0192965357819"}
		table_data.append(row_data1)
		row_data2 = {"name":"龙塘新村西(公交站)","distance":"1.8","lng":"113.1051356943977","lat":"23.0191965357819"}
		table_data.append(row_data2)
		row_data3 = {"name": "红山地铁站(公交站)", "distance":"2.8","lng": "113.1051356943977", "lat": "23.0191965357819"}
		table_data.append(row_data3)
		row_data4 = {"name": "深业泰富广场(公交站)", "distance":"3.8","lng": "113.1051356943977", "lat": "23.0191965357819"}
		table_data.append(row_data4)

		# 表格样式 --------------------
		rows = len(table_data)+1
		cols = 3
		top = Cm(2)
		left = Cm(16)  # Inches(2.0)
		width = Cm(10)  # Inches(6.0)
		height = Cm(2)  # Inches(0.8)

		# 添加表格到幻灯片 --------------------
		table = slide20.shapes.add_table(rows, cols, left, top, width, height).table
		# 设置单元格宽度
		table.columns[0].width = Cm(2)  # Inches(2.0)
		table.columns[1].width = Cm(4)
		table.columns[2].width = Cm(3)

		# 插入数据到table
		for row_index in range(len(table.rows)):
			if row_index == 0:
				table.rows[row_index].cells[0].text="编号"
				table.rows[row_index].cells[1].text = "名称"
				table.rows[row_index].cells[2].text = "距离(km)"
				continue
			for cell_index in range(len(table.rows[row_index].cells)):
				if cell_index == 0:
					table.rows[row_index].cells[cell_index].text = str(row_index)
				if cell_index == 1:
					table.rows[row_index].cells[cell_index].text = table_data[row_index-1].get("name")
				if cell_index == 2:
					table.rows[row_index].cells[cell_index].text = table_data[row_index-1].get("distance")

		# 处理table的样式、字体等
		for cell in diyun_ppt_page_handler.iter_cells(table):
					for paragraph in cell.text_frame.paragraphs:
						for run in paragraph.runs:
							run.font.size = Pt(12)
							run.font.color.rgb = RGBColor(0, 0, 0)  # 黑色字体

	@staticmethod
	def handler_page_21(slide21):
		table_data = []
		row_data1 = {"name":"和平实验小学","distance":"0.8","lng":"113.1051556943977","lat":"23.0192965357819"}
		table_data.append(row_data1)
		row_data2 = {"name":"民治小学","distance":"1.8","lng":"113.1051356943977","lat":"23.0191965357819"}
		table_data.append(row_data2)
		row_data3 = {"name": "深圳高级中学(集团)北校区", "distance":"2.8","lng": "113.1051356943977", "lat": "23.0191965357819"}
		table_data.append(row_data3)
		row_data4 = {"name": "希望中学", "distance":"3.8","lng": "113.1051356943977", "lat": "23.0191965357819"}
		table_data.append(row_data4)

		# 表格样式 --------------------
		rows = len(table_data)+1
		cols = 3
		top = Cm(2)
		left = Cm(16)  # Inches(2.0)
		width = Cm(10)  # Inches(6.0)
		height = Cm(2)  # Inches(0.8)

		# 添加表格到幻灯片 --------------------
		table = slide21.shapes.add_table(rows, cols, left, top, width, height).table
		# 设置单元格宽度
		table.columns[0].width = Cm(2)  # Inches(2.0)
		table.columns[1].width = Cm(4)
		table.columns[2].width = Cm(3)

		# 插入数据到table
		for row_index in range(len(table.rows)):
			if row_index == 0:
				table.rows[row_index].cells[0].text="编号"
				table.rows[row_index].cells[1].text = "名称"
				table.rows[row_index].cells[2].text = "距离(km)"
				continue
			for cell_index in range(len(table.rows[row_index].cells)):
				if cell_index == 0:
					table.rows[row_index].cells[cell_index].text = str(row_index)
				if cell_index == 1:
					table.rows[row_index].cells[cell_index].text = table_data[row_index-1].get("name")
				if cell_index == 2:
					table.rows[row_index].cells[cell_index].text = table_data[row_index-1].get("distance")

		# 处理table的样式、字体等
		for cell in diyun_ppt_page_handler.iter_cells(table):
					for paragraph in cell.text_frame.paragraphs:
						for run in paragraph.runs:
							run.font.size = Pt(12)
							run.font.color.rgb = RGBColor(0, 0, 0)  # 黑色字体


	@staticmethod
	def handler_page_22(slide22):
		table_data = []
		row_data1 = {"name":"深圳市第一人民医院","distance":"0.8","lng":"113.1051556943977","lat":"23.0192965357819"}
		table_data.append(row_data1)
		row_data2 = {"name":"龙岗人民医院","distance":"1.8","lng":"113.1051356943977","lat":"23.0191965357819"}
		table_data.append(row_data2)
		row_data3 = {"name": "香港中文大学附属第一人民医院", "distance":"2.8","lng": "113.1051356943977", "lat": "23.0191965357819"}
		table_data.append(row_data3)
		row_data4 = {"name": "南方科技大学附属第一人民医院", "distance":"3.8","lng": "113.1051356943977", "lat": "23.0191965357819"}
		table_data.append(row_data4)

		# 表格样式 --------------------
		rows = len(table_data)+1
		cols = 3
		top = Cm(2)
		left = Cm(16)  # Inches(2.0)
		width = Cm(10)  # Inches(6.0)
		height = Cm(2)  # Inches(0.8)

		# 添加表格到幻灯片 --------------------
		table = slide22.shapes.add_table(rows, cols, left, top, width, height).table
		# 设置单元格宽度
		table.columns[0].width = Cm(2)  # Inches(2.0)
		table.columns[1].width = Cm(4)
		table.columns[2].width = Cm(3)

		# 插入数据到table
		for row_index in range(len(table.rows)):
			if row_index == 0:
				table.rows[row_index].cells[0].text="编号"
				table.rows[row_index].cells[1].text = "名称"
				table.rows[row_index].cells[2].text = "距离(km)"
				continue
			for cell_index in range(len(table.rows[row_index].cells)):
				if cell_index == 0:
					table.rows[row_index].cells[cell_index].text = str(row_index)
				if cell_index == 1:
					table.rows[row_index].cells[cell_index].text = table_data[row_index-1].get("name")
				if cell_index == 2:
					table.rows[row_index].cells[cell_index].text = table_data[row_index-1].get("distance")

		# 处理table的样式、字体等
		for cell in diyun_ppt_page_handler.iter_cells(table):
					for paragraph in cell.text_frame.paragraphs:
						for run in paragraph.runs:
							run.font.size = Pt(12)
							run.font.color.rgb = RGBColor(0, 0, 0)  # 黑色字体

	@staticmethod
	def handler_page_23(slide23):
		table_data = []
		row_data1 = {"name": "九方购物中心(人民路)", "distance": "0.8", "lng": "113.1051556943977", "lat": "23.0192965357819"}
		table_data.append(row_data1)
		row_data2 = {"name": "中海环宇新天地", "distance": "1.8", "lng": "113.1051356943977", "lat": "23.0191965357819"}
		table_data.append(row_data2)
		row_data3 = {"name": "U·ONE优城", "distance": "2.8", "lng": "113.1051356943977", "lat": "23.0191965357819"}
		table_data.append(row_data3)
		row_data4 = {"name": "汇龙时代广场", "distance": "3.8", "lng": "113.1051356943977", "lat": "23.0191965357819"}
		table_data.append(row_data4)

		# 表格样式 --------------------
		rows = len(table_data) + 1
		cols = 3
		top = Cm(2)
		left = Cm(16)  # Inches(2.0)
		width = Cm(10)  # Inches(6.0)
		height = Cm(2)  # Inches(0.8)

		# 添加表格到幻灯片 --------------------
		table = slide23.shapes.add_table(rows, cols, left, top, width, height).table
		# 设置单元格宽度
		table.columns[0].width = Cm(2)  # Inches(2.0)
		table.columns[1].width = Cm(4)
		table.columns[2].width = Cm(3)

		# 插入数据到table
		for row_index in range(len(table.rows)):
			if row_index == 0:
				table.rows[row_index].cells[0].text = "编号"
				table.rows[row_index].cells[1].text = "名称"
				table.rows[row_index].cells[2].text = "距离(km)"
				continue
			for cell_index in range(len(table.rows[row_index].cells)):
				if cell_index == 0:
					table.rows[row_index].cells[cell_index].text = str(row_index)
				if cell_index == 1:
					table.rows[row_index].cells[cell_index].text = table_data[row_index - 1].get("name")
				if cell_index == 2:
					table.rows[row_index].cells[cell_index].text = table_data[row_index - 1].get("distance")

		# 处理table的样式、字体等
		for cell in diyun_ppt_page_handler.iter_cells(table):
			for paragraph in cell.text_frame.paragraphs:
				for run in paragraph.runs:
					run.font.size = Pt(12)
					run.font.color.rgb = RGBColor(0, 0, 0)  # 黑色字体

	@staticmethod
	def handler_page_24(slide24):
		table_data = []
		row_data1 = {"name": "深圳市福田区政府", "distance": "0.8", "lng": "113.1051556943977", "lat": "23.0192965357819"}
		table_data.append(row_data1)
		row_data2 = {"name": "深圳市南山区政府", "distance": "1.8", "lng": "113.1051356943977", "lat": "23.0191965357819"}
		table_data.append(row_data2)
		row_data3 = {"name": "深圳市宝安区政府", "distance": "2.8", "lng": "113.1051356943977", "lat": "23.0191965357819"}
		table_data.append(row_data3)
		row_data4 = {"name": "宝山街道居委会", "distance": "3.8", "lng": "113.1051356943977", "lat": "23.0191965357819"}
		table_data.append(row_data4)

		# 表格样式 --------------------
		rows = len(table_data) + 1
		cols = 3
		top = Cm(2)
		left = Cm(16)  # Inches(2.0)
		width = Cm(10)  # Inches(6.0)
		height = Cm(2)  # Inches(0.8)

		# 添加表格到幻灯片 --------------------
		table = slide24.shapes.add_table(rows, cols, left, top, width, height).table
		# 设置单元格宽度
		table.columns[0].width = Cm(2)  # Inches(2.0)
		table.columns[1].width = Cm(4)
		table.columns[2].width = Cm(3)

		# 插入数据到table
		for row_index in range(len(table.rows)):
			if row_index == 0:
				table.rows[row_index].cells[0].text = "编号"
				table.rows[row_index].cells[1].text = "名称"
				table.rows[row_index].cells[2].text = "距离(km)"
				continue
			for cell_index in range(len(table.rows[row_index].cells)):
				if cell_index == 0:
					table.rows[row_index].cells[cell_index].text = str(row_index)
				if cell_index == 1:
					table.rows[row_index].cells[cell_index].text = table_data[row_index - 1].get("name")
				if cell_index == 2:
					table.rows[row_index].cells[cell_index].text = table_data[row_index - 1].get("distance")

		# 处理table的样式、字体等
		for cell in diyun_ppt_page_handler.iter_cells(table):
			for paragraph in cell.text_frame.paragraphs:
				for run in paragraph.runs:
					run.font.size = Pt(12)
					run.font.color.rgb = RGBColor(0, 0, 0)  # 黑色字体

	@staticmethod
	def handler_test_in_last_page(presentation):
		# 使用默认样式，追加一个空白的ppt,并插入柱状图和表格
		slide = presentation.slides.add_slide(presentation.slide_layouts[6])
		shapes = slide.shapes

		# 定义图表数据 ---------------------
		# 堆积柱形图：定义图表数据 ---------------------
		chart_data = ChartData()
		chart_data.categories = ['2014', '2015', '2016', '2017', '2018']
		production1data = (19.2, 21.4, 16.7,13,15)
		production2data = (9.2, 2.4, 6.7,12,13.5)
		production3data = (9.2, 2.4, 6.7,14.1,5.2)
		productionDataMap ={}
		productionDataMap["production1data"] = production1data
		productionDataMap["production2data"] = production2data
		productionDataMap["production3data"] = production3data

		chart_data.add_series('第一产业', production1data)
		chart_data.add_series('第二产业', production2data)
		chart_data.add_series('第三产业', production3data)


		# 堆积柱形图：将图表添加到幻灯片 --------------------
		left, top, width, height = Inches(1), Inches(1.3), Inches(8), Inches(3.5)
		graphic_frame = slide.shapes.add_chart(
			XL_CHART_TYPE.COLUMN_STACKED, left, top, width, height, chart_data
		)
		chart = graphic_frame.chart
		chart.has_legend = True
		chart.legend.position = XL_LEGEND_POSITION.BOTTOM
		chart.legend.include_in_layout = False
		graphic_frame.has_title = True
		graphic_frame.chart.chart_title.text_frame.text = "深圳市产业结构"


		# 定制柱状图颜色
		colors = {}
		colors['第一产业'] = 'FFD700'
		colors['第二产业'] = '32CD32'
		colors['第三产业'] = '1E90FF'
		for series in chart.series:
			if series.name in colors:
				fill = series.format.fill  # fill the legend as well
				fill.solid()
				fill.fore_color.rgb = RGBColor.from_string(colors[series.name])


		# 把数据填充到table中
		graphicFrames = slide.shapes
		for graphicFrame in graphicFrames:
			if type(graphicFrame) != GraphicFrame or not graphicFrame.has_table:
				continue
			for row_index in range(len(graphicFrame.table.rows)):
				if row_index == 0:
					continue
				row_data = productionDataMap.get("production{}data".format(row_index))
				for column_index in range(len(graphicFrame.table.rows[row_index].cells)):
					if column_index == 0:
						continue
					# 插入数据到 每个单元格
					graphicFrame.table.rows[row_index].cells[column_index].text = str(row_data[column_index-1])

					# 处理 每个单元格 的样式
					for paragraph in graphicFrame.table.rows[row_index].cells[column_index].text_frame.paragraphs:
						for run in paragraph.runs:
							run.font.size = Pt(12)
							run.font.color.rgb = RGBColor(0, 0, 0)  # 黑色字体

		# 定义表格数据 ------
		name_objects = ["object1", "object2", "object3"]
		name_AIs = ["AI1", "AI2", "AI3"]
		val_AI1 = (19.2, 21.4, 16.7)
		val_AI2 = (22.3, 28.6, 15.2)
		val_AI3 = (20.4, 26.3, 14.2)
		val_AIs = [val_AI1, val_AI2, val_AI3]
		# 表格样式 --------------------
		rows = 4
		cols = 4
		top = Cm(12.5)
		left = Cm(0.5)  # Inches(2.0)
		width = Cm(3)  # Inches(6.0)
		height = Cm(0.8)  # Inches(0.8)

		# 添加表格到幻灯片 --------------------
		table = shapes.add_table(rows, cols, left, top, width, height).table

		# 设置单元格宽度
		table.columns[0].width = Cm(3)  # Inches(2.0)
		table.columns[1].width = Cm(3)
		table.columns[2].width = Cm(3)
		table.columns[3].width = Cm(3)

		# 设置标题行
		table.cell(0, 1).text = name_objects[0]
		table.cell(0, 2).text = name_objects[1]
		table.cell(0, 3).text = name_objects[2]

		# 填充数据
		table.cell(1, 0).text = name_AIs[0]
		table.cell(1, 1).text = str(val_AI1[0])
		table.cell(1, 2).text = str(val_AI1[1])
		table.cell(1, 3).text = str(val_AI1[2])

		table.cell(2, 0).text = name_AIs[1]
		table.cell(2, 1).text = str(val_AI2[0])
		table.cell(2, 2).text = str(val_AI2[1])
		table.cell(2, 3).text = str(val_AI2[2])

		table.cell(3, 0).text = name_AIs[2]
		table.cell(3, 1).text = str(val_AI3[0])
		table.cell(3, 2).text = str(val_AI3[1])
		table.cell(3, 3).text = str(val_AI3[2])


