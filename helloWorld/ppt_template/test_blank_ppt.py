from pptx import Presentation

# 创建幻灯片 ------
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)

title = slide.shapes.title
subtitle = slide.placeholders[1]

# 设置标题和副标题
title.text = "Hello, World!"
subtitle.text = "pip install python-pptx"

prs.save("test_blank_ppt.pptx")