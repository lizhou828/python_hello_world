import time

from pptx import Presentation

# 创建幻灯片 ------
from helloWorld.ppt_template.diyun_ppt_page_handler import diyun_ppt_page_handler

start_time = time.time()

prs = Presentation("diyun_land_report_template.pptx") #当前文件目录下打开已有的PPT

#找到第1页PPT，文字替换处理
page1_start_time = time.time()
diyun_ppt_page_handler.handler_page_1(prs.slides[0])
page1_seconds_spend = int(time.time() - page1_start_time)
print("PPT第1页处理完成，耗时：%d 秒" % page1_seconds_spend)

#处理第4页
page4_start_time = time.time()
diyun_ppt_page_handler.handler_page_4(prs.slides[3])
page4_seconds_spend = int(time.time() - page4_start_time)
print("PPT第4页处理完成，耗时：%d 秒" % page4_seconds_spend)

#处理第5页，table中直接插入数据
diyun_ppt_page_handler.handler_page_5(prs.slides[4])

#处理第12页，插入百度地图的截图
page12_start_time = time.time()
diyun_ppt_page_handler.handler_page_12(prs.slides[11])
page12_seconds_spend = int(time.time() - page12_start_time)
print("PPT第12页，插入地图的截图完成，耗时：%d 秒" % page12_seconds_spend)

#处理第14页，插入3个table
diyun_ppt_page_handler.handler_page_14_new(prs.slides[13])


#处理第18页，table中动态插入多行
diyun_ppt_page_handler.handler_page_18(prs.slides[17])

#处理第20-24页，插入table及其数据
diyun_ppt_page_handler.handler_page_20(prs.slides[19])
diyun_ppt_page_handler.handler_page_21(prs.slides[20])
diyun_ppt_page_handler.handler_page_22(prs.slides[21])
diyun_ppt_page_handler.handler_page_23(prs.slides[22])
diyun_ppt_page_handler.handler_page_24(prs.slides[23])



#追加一页默认样式的空白页，测试插入图标和表格
lastpage_start_time = time.time()
diyun_ppt_page_handler.handler_test_in_last_page(prs)
lastpage_seconds_spend = int(time.time() - lastpage_start_time)
print("PPT最后一页的图标、表格生成完成，耗时：%d 秒" % lastpage_seconds_spend)


prs.save("diyun_land_report_result.pptx")

seconds_spend = int(time.time() - start_time)

print("生成PPT完成，总耗时：%d 秒" % seconds_spend)
